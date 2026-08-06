# -*- coding: utf-8 -*-
"""Dek 16 slide untuk reports-simple. Bahasa Indonesia, gaya lisan seminar.

Tata letak dikunci lewat konstanta baris di bawah supaya tidak ada kotak yang
bertumpuk. Verifikasi geometri ada di verify_deck.py.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = r"D:\Work\Assisten-Dosen\Research-Pipeline\output\presentation\reports-simple-id.pptx"

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
TEAL  = RGBColor(0x2E, 0x8B, 0x8B)
RED   = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xD6, 0x89, 0x10)
INK   = RGBColor(0x22, 0x2B, 0x35)
MUTED = RGBColor(0x5A, 0x66, 0x72)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
CREAM = RGBColor(0xFD, 0xF6, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINEG = RGBColor(0xD8, 0xDF, 0xE6)

# Kisi vertikal bersama (inci)
M       = 0.80    # margin kiri
CW      = 11.75   # lebar kolom isi
KICKER  = 0.34
TITLE_Y = 0.70
BODY_Y  = 1.66    # batas atas isi
FOOT_Y  = 6.98

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

RECT, ROUND = 1, 5


def slide():
    return prs.slides.add_slide(BLANK)


def tbox(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    return tf


def par(tf, text, size=18, bold=False, color=INK, after=0, before=0,
        align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return p


def rect(s, l, t, w, h, fill, line=None, kind=RECT):
    sh = s.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    return sh


def header(s, title, kicker):
    rect(s, 0, 0, 13.333, 0.09, NAVY)
    par(tbox(s, M, KICKER, CW, 0.30), kicker.upper(), 12.5, True, TEAL, first=True)
    par(tbox(s, M, TITLE_Y, CW, 0.72), title, 29, True, NAVY, first=True)


def foot(s, n):
    par(tbox(s, 12.15, FOOT_Y, 0.85, 0.30), str(n), 11.5, False, MUTED,
        align=PP_ALIGN.RIGHT, first=True)


def callout(s, lines, l, t, w, h, fill=LIGHT, edge=NAVY, size=16.5,
            bold_first=True, color=INK, align=PP_ALIGN.LEFT):
    sh = rect(s, l, t, w, h, fill, edge, kind=ROUND)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.26)
    tf.margin_top = tf.margin_bottom = Inches(0.10)
    for i, line in enumerate(lines):
        par(tf, line, size, bold_first and i == 0, color, 3, align=align,
            first=(i == 0))
    return sh


def pic(s, name, l, t, w):
    return s.shapes.add_picture(os.path.join(HERE, name), Inches(l), Inches(t),
                                width=Inches(w))


def lead(s, text, size=17, t=BODY_Y, h=0.64):
    par(tbox(s, M, t, CW, h), text, size, False, MUTED, first=True)


def blocks(s, items, top, pitch, head_size=18, sub_size=15):
    """Blok kepala tebal + satu kalimat penjelas."""
    for i, (head, sub) in enumerate(items):
        y = top + i * pitch
        rect(s, M, y + 0.06, 0.05, 0.30, TEAL)
        par(tbox(s, M + 0.22, y, CW - 0.22, 0.36), head, head_size, True, NAVY,
            first=True)
        par(tbox(s, M + 0.22, y + 0.38, CW - 0.22, pitch - 0.44), sub, sub_size,
            False, MUTED, first=True)


def table(s, cols, rows, top, row_h, pitch, sizes, bolds, colors, header_row=None):
    """cols = daftar (offset_kiri, lebar) dalam inci, relatif terhadap M."""
    y = top
    if header_row:
        rect(s, M, y, CW, row_h - 0.06, NAVY)
        for (dx, w), txt in zip(cols, header_row):
            par(tbox(s, M + dx + 0.12, y + 0.09, w - 0.20, row_h - 0.24), txt,
                13, True, WHITE, first=True)
        y += pitch
    for i, r in enumerate(rows):
        rect(s, M, y, CW, row_h, LIGHT if i % 2 == 0 else WHITE)
        for (dx, w), txt, sz, bd, cl in zip(cols, r, sizes, bolds, colors):
            par(tbox(s, M + dx + 0.12, y + 0.10, w - 0.20, row_h - 0.20), txt,
                sz, bd, cl, first=True)
        y += pitch
    return y


# ===================================================== 1. Judul
s = slide()
rect(s, 0, 0, 13.333, 3.50, NAVY)
par(tbox(s, M + 0.10, 0.92, CW, 0.32),
    "LAPORAN EKSPERIMEN \u00b7 DETEKSI TANDAN BUAH SEGAR KELAPA SAWIT",
    13.5, True, RGBColor(0x8F, 0xC7, 0xC7), first=True)
par(tbox(s, M + 0.10, 1.36, CW, 0.62), "Menguji Kanal Keempat:", 34, True,
    WHITE, first=True)
par(tbox(s, M + 0.10, 1.98, CW, 0.62), "Apakah Depth Membantu Deteksi?", 34,
    True, WHITE, first=True)
par(tbox(s, M + 0.10, 2.72, CW, 0.42),
    "Ringkasan 32 percobaan dan 9 pemeriksaan bukti", 18, False,
    RGBColor(0xC7, 0xD6, 0xE2), first=True)

par(tbox(s, M + 0.10, 3.86, CW, 0.42),
    "Muhammad Zainal Muttaqin  \u00b7  Fatma Indriani", 20, True, INK, first=True)
par(tbox(s, M + 0.10, 4.30, CW, 0.34), "Naskah pra-cetak \u00b7 4 Agustus 2026",
    15, False, MUTED, first=True)

callout(s, [
    "Jawaban singkatnya: pada kondisi yang diuji, kanal depth belum menaikkan akurasi.",
    "Datanya juga belum cukup untuk menyimpulkan bahwa gagasan depth itu keliru.",
], M, 5.05, CW, 1.20, fill=CREAM, edge=AMBER, size=16.5)

# ===================================================== 2. Masalah
s = slide()
header(s, "Masalah di lapangan", "latar")
blocks(s, [
    ("Kamera biasa merekam warna saja.",
     "Satu foto terdiri atas tiga lapis, yaitu merah, hijau, dan biru."),
    ("Dua tandan dapat berwarna mirip dan saling menutupi.",
     "Detektor lalu menyatukan keduanya, atau melewatkan salah satunya."),
    ("Tandan mentah kelas B4 paling sering luput.",
     "Warnanya kehijauan gelap dan menyatu dengan pelepah."),
], top=BODY_Y, pitch=0.95)

callout(s, [
    "Kamera depth menambah satu lapis lagi, yaitu jarak ke setiap piksel.",
    "Dua tandan yang warnanya sama tetap berada pada jarak berbeda. Lapis jarak "
    "dapat memisahkan keduanya.",
], M, 4.72, CW, 1.32, size=17.5)
par(tbox(s, M, 6.32, 10.9, 0.42),
    "Kamera yang dipakai adalah Orbbec Gemini. Keluarannya satu nilai jarak per piksel.",
    14, False, MUTED, first=True)
foot(s, 2)

# ===================================================== 3. Pertanyaan
s = slide()
header(s, "Satu pertanyaan yang ingin dijawab", "rumusan masalah")
lead(s, "Lapis jarak ditambahkan sebagai masukan keempat. Akurasi deteksi hanya "
        "punya tiga kemungkinan arah.")

opts = [("NAIK", "Depth membawa informasi\nyang berguna.", TEAL),
        ("TETAP", "Depth tidak menambah\napa pun.", RGBColor(0x8A, 0x94, 0x9E)),
        ("TURUN", "Depth menambah derau\nke dalam masukan.", RED)]
for i, (head, sub, col) in enumerate(opts):
    l = M + i * 3.95
    sh = rect(s, l, 2.48, 3.55, 1.72, WHITE, col, kind=ROUND)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    par(tf, head, 26, True, col, 5, align=PP_ALIGN.CENTER, first=True)
    for line in sub.split("\n"):
        par(tf, line, 15, False, INK, 1, align=PP_ALIGN.CENTER)

callout(s, [
    "Kami menjawabnya dengan melatih detektor dengan dan tanpa lapis jarak, lalu "
    "membandingkan skornya pada uji yang sama persis.",
], M, 4.62, CW, 0.92, size=17.5)
par(tbox(s, M, 5.78, 10.9, 0.80),
    "Ukurannya adalah mAP50, yaitu rata-rata ketepatan deteksi pada seluruh kelas kematangan. "
    "Nilainya berkisar 0 sampai 1.",
    14, False, MUTED, first=True)
foot(s, 3)

# ===================================================== 4. Cara kerja
s = slide()
header(s, "Bagaimana penelitian ini dijalankan", "metode")
lead(s, "Setiap percobaan memegang satu hipotesis yang dapat dipatahkan, dan "
        "hasilnya dicatat pada log bernomor.")

table(s,
      cols=[(0.0, 3.60), (3.60, 2.05), (5.65, 6.10)],
      rows=[
          ("Satu hipotesis per percobaan", "32 percobaan",
           "Setiap percobaan menjawab satu pertanyaan kecil."),
          ("Pemeriksaan sebelum simpulan", "9 pemeriksaan",
           "Setiap lubang pada bukti ditutup lebih dahulu."),
          ("Hasil negatif dicatat", "wajib",
           "Percobaan yang gagal ditulis sama seperti yang berhasil."),
          ("Hasil cacat ditarik", "berlaku",
           "Angka positif pertama dibatalkan setelah galat ditemukan."),
      ],
      top=2.42, row_h=0.82, pitch=0.90,
      sizes=(16.5, 14.5, 14.5), bolds=(True, True, False),
      colors=(NAVY, TEAL, INK))

par(tbox(s, M, 6.20, 10.9, 0.62),
    "Percobaannya dijalankan secara otomatis. Penomoran menjaga arah kerjanya "
    "supaya tidak melebar.", 14, False, MUTED, first=True)
foot(s, 4)

# ===================================================== 5. Peta percobaan
s = slide()
header(s, "Delapan percobaan yang paling menentukan", "peta percobaan")
lead(s, "Delapan percobaan ini membawa hasil utamanya. Sisanya menutup "
        "penjelasan alternatif.", size=15.5, t=1.58, h=0.40)

table(s,
      cols=[(0.0, 3.30), (3.30, 4.90), (8.20, 3.55)],
      header_row=("Nama percobaan", "Pertanyaannya", "Hasilnya"),
      rows=[
          ("Uji konsistensi label",
           "Apakah label kematangan konsisten?", "Tidak. Labelnya bersih."),
          ("Uji depth tiruan",
           "Dapatkah jarak ditebak dari foto biasa?", "Gagal. Kontrasnya lemah."),
          ("Uji penyebab B4 sulit",
           "Mengapa tandan mentah sering luput?", "Karena kontras, bukan desakan."),
          ("Adu detektor warna",
           "Detektor RGB mana yang jadi pembanding?", "RF-DETR-L, mAP50 0,6038."),
          ("Audit berkas sensor",
           "Apakah berkas depth sesuai metadatanya?", "Tidak. Lima galat ditemukan."),
          ("Uji kesamaan alat ukur",
           "Apakah kedua sisi diukur dengan cara sama?", "Belum. Alat ukur disatukan."),
          ("Uji model lebih besar",
           "Apakah model besar menyelamatkan depth?", "Tidak terbukti."),
          ("Uji tiga titik fusi",
           "Di lapisan mana depth sebaiknya digabung?", "Ketiganya sama saja."),
      ],
      top=2.12, row_h=0.51, pitch=0.545,
      sizes=(14, 12.5, 12.5), bolds=(True, False, True),
      colors=(NAVY, INK, TEAL))
foot(s, 5)

# ===================================================== 6. Detektor warna
s = slide()
header(s, "Detektor warna terbaik sebagai titik acuan", "hasil rgb")
pic(s, "chart_rgb.png", 0.87, 1.52, 11.60)   # tinggi 4,58
callout(s, [
    "RF-DETR-L mencapai mAP50 0,6038 pada dataset warna yang lama dan melewati "
    "sasaran 0,60. Angka ini memakai dataset lain, sehingga tidak dapat "
    "dibandingkan dengan angka pada slide depth.",
], M, 6.16, CW, 0.72, size=15.5)
foot(s, 6)

# ===================================================== 7. Audit
s = slide()
header(s, "Alat ukurnya diperiksa lebih dahulu", "audit")
lead(s, "Hasil negatif hanya berguna apabila cara mengukurnya benar. "
        "Pemeriksaan menemukan lima galat, dan kelimanya sudah diperbaiki.")

defects = [
    ("Foto dan jarak tidak sejajar",
     "Melesetnya 29 piksel. Diperbaiki dengan reproyeksi penuh."),
    ("Kalibrasi dua kamera disamakan",
     "Kalibrasi kini dibaca satu per satu dari tiap berkas."),
    ("Rentang jarak salah",
     "Rentang 0,3\u20138 m diganti 0,8\u201315 m sesuai sensornya."),
    ("Data bocor antar-pembagian",
     "Hasil yang terpengaruh sudah ditarik."),
    ("Alat ukur kedua sisi berbeda",
     "Cukup untuk membalik tanda hasil. Alat ukur kini satu."),
]
y0 = 2.36
for i, (a, b) in enumerate(defects):
    y = y0 + i * 0.68
    rect(s, M, y, 0.06, 0.56, RED)
    par(tbox(s, M + 0.24, y + 0.02, 4.60, 0.34),
        f"{i + 1}.  {a}", 16, True, INK, first=True)
    par(tbox(s, M + 5.00, y + 0.05, 6.55, 0.44), b, 14.5, False, MUTED, first=True)

callout(s, [
    "Hasil positif yang muncul pertama ternyata berasal dari alur yang cacat ini, "
    "dan sudah dibatalkan. Seluruh angka pada slide berikutnya berasal dari alur "
    "yang sudah diperbaiki.",
], M, 5.96, CW, 0.88, fill=CREAM, edge=AMBER, size=15.5)
foot(s, 7)

# ===================================================== 8. Tiga titik fusi
s = slide()
header(s, "Tiga cara menggabungkan jarak ke dalam jaringan", "rancangan")
pic(s, "fig_hipotesis.png", 0.80, 1.52, 7.05)   # tinggi 4,70

kunci = [
    ("Early fusion (fusi awal)",
     "RGB dan lapis jarak disatukan sejak masukan."),
    ("Middle fusion (fusi menengah)",
     "Keduanya diproses sendiri, lalu digabung di tengah."),
    ("Late fusion (fusi akhir)",
     "Dua jaringan penuh, digabung menjelang keluaran."),
    ("Noise control (kendali derau)",
     "Jarak mentah disaring, lubangnya ditambal, lalu dihaluskan."),
]
par(tbox(s, 8.20, 1.52, 4.35, 0.36), "Bacaan singkat", 17.5, True, NAVY, first=True)
yk = 2.00
for a, b in kunci:
    rect(s, 8.20, yk + 0.05, 0.05, 0.26, TEAL)
    par(tbox(s, 8.40, yk, 4.15, 0.32), a, 14, True, INK, first=True)
    par(tbox(s, 8.40, yk + 0.34, 4.15, 0.56), b, 12.5, False, MUTED, first=True)
    yk += 0.94
callout(s, [
    "Depth adalah hipotesis, bukan jaminan kenaikan.",
    "Manfaatnya bergantung pada mutu sensor dan rancangan fusi.",
], 8.20, 5.84, 4.35, 1.14, fill=CREAM, edge=RED, size=12.5)
foot(s, 8)

# ===================================================== 9. Hasil utama
s = slide()
header(s, "Hasil utama pada ketiga titik penggabungan", "hasil")
pic(s, "chart_fusi.png", 0.77, 1.52, 11.80)   # tinggi 4,52
callout(s, [
    "Batangnya naik-turun tanpa pola. Tidak satu kelompok pun keluar dari pita "
    "ragam yang muncul hanya karena pengacakan awal model.",
], M, 6.16, CW, 0.72, size=15.5)
foot(s, 9)

# ===================================================== 10. Derau acak
s = slide()
header(s, "Kanal berisi angka acak sebagai pembanding", "hasil")
lead(s, "Lapis jarak diganti bilangan acak, yaitu kanal yang pasti tidak memuat "
        "informasi bentuk.", size=16, t=1.58, h=0.44)
pic(s, "chart_derau.png", 1.15, 2.10, 11.00)  # tinggi 3,40
callout(s, [
    "Kanal acak menggeser skor sebesar lapis jarak. Pada kondisi ini detektor "
    "belum memanfaatkan isi lapis jarak.",
], M, 5.62, CW, 0.84, edge=RED, size=16)
par(tbox(s, M, 6.56, 10.9, 0.42),
    "Ini bukan uji ekuivalensi. Pernyataan yang sah adalah \"belum terlihat "
    "manfaatnya\", bukan \"depth sama dengan derau\".", 13.5, False, MUTED, first=True)
foot(s, 10)

# ===================================================== 11. Batas data
s = slide()
header(s, "Datanya belum memadai untuk uji ini", "batas bukti")
pic(s, "chart_dataset.png", 1.27, 1.46, 10.80)  # tinggi 4,30
callout(s, [
    "Dataset depth memuat sepertiga jumlah pohon, sepertiga kepadatan tandan, dan "
    "sebaran kelas yang terbalik. Angka mAP di sini tidak dapat dibandingkan "
    "langsung dengan dataset lama. Hasil negatif ini menguji satu penerapan depth, "
    "bukan gagasan depth itu sendiri.",
], M, 5.86, CW, 1.06, fill=CREAM, edge=AMBER, size=15.5)
foot(s, 11)

# ===================================================== 12. Contoh citra lama
s = slide()
header(s, "Tandan pada dataset warna yang lama", "contoh citra")
pic(s, "panel_mvc.png", 1.42, 1.46, 10.50)
callout(s, [
    "Batas tandan terlihat tegas. Buah matang berwarna jingga-merah pekat, dan "
    "tandan mentah tetap terbaca meski gelap.",
], M, 6.04, CW, 0.86, size=15.5)
foot(s, 12)

# ===================================================== 13. Contoh citra depth
s = slide()
header(s, "Tandan pada dataset depth", "contoh citra")
pic(s, "panel_depth.png", 1.42, 1.46, 10.50)
callout(s, [
    "Potongan tandan di sini 3,4 kali lebih kabur dan 1,7 kali lebih pucat. "
    "Piksel tandannya justru lebih banyak, jadi yang kurang adalah mutu citra, "
    "bukan ukuran objek.",
], M, 6.04, CW, 0.86, fill=CREAM, edge=AMBER, size=15.5)
foot(s, 13)

# ===================================================== 14. Implikasi
s = slide()
header(s, "Arti hasil ini bagi sistem di lapangan", "implikasi")
blocks(s, [
    ("Pakai jalur warna untuk sekarang.",
     "RF-DETR-L menjadi pembanding terbaik yang teramati, dengan mAP50 0,6038. Kecepatannya 8,5 FPS "
     "pada GPU L4, sehingga masih perlu optimasi untuk waktu nyata."),
    ("Jangan membeli perangkat depth semata demi akurasi.",
     "Pengukuran yang ada belum mendukung keputusan itu."),
    ("Perangkat depth yang sudah ada tetap layak dipakai merekam.",
     "Simpan hasilnya sebagai cabang data terpisah, dengan penyaringan mutu per berkas."),
    ("Kekeliruan B2 lawan B3 perlu penanganan lain.",
     "Gunakan loss ordinal atau kepala regresi. Kekeliruannya berurutan, bukan acak."),
], top=BODY_Y, pitch=1.22, head_size=17.5, sub_size=14.5)
foot(s, 14)

# ===================================================== 15. Batas
s = slide()
header(s, "Batas jawaban ini dan pekerjaan yang tersisa", "batas hasil")

par(tbox(s, M, BODY_Y, 5.55, 0.38), "Kondisi yang diuji", 18.5, True, NAVY,
    first=True)
tf = tbox(s, M, BODY_Y + 0.46, 5.55, 3.55)
for i, t in enumerate([
    "Modelnya hanya YOLO26n, 640 piksel, 150 epoch.",
    "Datanya satu pembagian dengan tiga pengacakan awal.",
    "Pelatihannya dari nol, tanpa bobot awal.",
    "Dataset depth memuat 352 pohon.",
    "Kameranya satu jenis dengan satu penyandian jarak.",
]):
    par(tf, "\u25B8  " + t, 15.5, False, INK, 11, first=(i == 0))

par(tbox(s, 7.00, BODY_Y, 5.55, 0.38), "Yang belum diuji", 18.5, True, RED,
    first=True)
tf = tbox(s, 7.00, BODY_Y + 0.46, 5.55, 3.55)
for i, t in enumerate([
    "Pengaruh ukuran model baru diuji satu kali.",
    "Lima dari 12 pelatihan tambahan belum selesai.",
    "Perbandingan berpasangannya belum dihitung.",
    "Detektor RGB terbaik belum diuji dengan depth.",
    "Kepadatan tandan yang setara belum pernah diuji.",
]):
    par(tf, "\u25B8  " + t, 15.5, False, INK, 11, first=(i == 0))

callout(s, [
    "Kesimpulan ini bersyarat. Mengubah satu saja kondisi di kolom kiri dapat "
    "mengubah jawabannya.",
], M, 5.92, CW, 0.80, size=16.5)
foot(s, 15)

# ===================================================== 16. Simpulan
s = slide()
rect(s, 0, 0, 13.333, 0.09, NAVY)
par(tbox(s, M, 0.46, CW, 0.66), "Simpulan", 31, True, NAVY, first=True)

cards = [
    ("1", "Manfaatnya belum terukur.",
     "Ketiga titik penggabungan bergerak lebih kecil daripada ragam pengacakan model "
     "itu sendiri. Seluruh 12 selang memuat nol.", TEAL),
    ("2", "Alat ukurnya sudah bersih.",
     "Lima galat diperbaiki lebih dahulu, dan hasil positif yang cacat sudah ditarik. "
     "Hasil negatif ini bukan akibat salah hitung.", NAVY),
    ("3", "Gagasannya belum tertutup.",
     "Datanya kecil dan jarang isinya. Yang terbukti satu penerapan depth belum "
     "berhasil, bukan bahwa depth tidak berguna.", AMBER),
]
y = 1.42
for num, head, body, col in cards:
    sh = rect(s, M, y, 0.72, 1.34, col)
    tfn = sh.text_frame
    tfn.vertical_anchor = MSO_ANCHOR.MIDDLE
    par(tfn, num, 28, True, WHITE, align=PP_ALIGN.CENTER, first=True)
    par(tbox(s, M + 0.95, y + 0.04, 10.60, 0.42), head, 20, True, INK, first=True)
    par(tbox(s, M + 0.95, y + 0.52, 10.60, 0.80), body, 15.5, False, MUTED,
        first=True)
    y += 1.56

callout(s, [
    "Langkah berikutnya menguji pengaruh ukuran model pada lebih dari satu "
    "pengacakan awal, lalu mengulang uji depth pada dataset yang kepadatan "
    "tandannya setara dataset lama.",
], M, 6.12, CW, 0.84, size=16.5)
foot(s, 16)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("tersimpan:", OUT)
