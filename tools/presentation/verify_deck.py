# -*- coding: utf-8 -*-
"""Periksa dek: kotak di luar kanvas, teks bertumpuk, dan taksiran luapan teks."""
import sys, math
from pptx import Presentation
from pptx.util import Emu

PATH = sys.argv[1] if len(sys.argv) > 1 else \
    r"D:\Work\Assisten-Dosen\Research-Pipeline\output\presentation\reports-simple-id.pptx"
EMU = 914400.0
prs = Presentation(PATH)
SW, SH = prs.slide_width / EMU, prs.slide_height / EMU

# Lebar rata-rata glyph Segoe UI, dalam satuan em. Diambil konservatif.
AVG_EM = 0.50
LINE_F = 1.22          # tinggi baris terhadap ukuran font
TOL = 0.02             # toleransi tumpang tindih, inci


def text_shapes(slide):
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if not sh.text_frame.text.strip():
            continue
        out.append(sh)
    return out


def need_height(sh):
    """Taksiran tinggi teks (inci) di dalam sebuah shape."""
    tf = sh.text_frame
    ml = (tf.margin_left or 0) / EMU
    mr = (tf.margin_right or 0) / EMU
    mt = (tf.margin_top or 0) / EMU
    mb = (tf.margin_bottom or 0) / EMU
    avail_w = sh.width / EMU - ml - mr
    total = mt + mb
    for p in tf.paragraphs:
        txt = "".join(r.text for r in p.runs)
        size = max([(r.font.size.pt if r.font.size else 18) for r in p.runs] or [18])
        if not txt:
            total += size / 72.0 * LINE_F
            continue
        char_w = size / 72.0 * AVG_EM
        cpl = max(1, int(avail_w / char_w))
        lines = max(1, math.ceil(len(txt) / cpl))
        total += lines * size / 72.0 * LINE_F
        total += (p.space_after.pt if p.space_after else 0) / 72.0
        total += (p.space_before.pt if p.space_before else 0) / 72.0
    return total


def rects_overlap(a, b):
    ax1, ay1 = a.left / EMU, a.top / EMU
    ax2, ay2 = ax1 + a.width / EMU, ay1 + a.height / EMU
    bx1, by1 = b.left / EMU, b.top / EMU
    bx2, by2 = bx1 + b.width / EMU, by1 + b.height / EMU
    ox = min(ax2, bx2) - max(ax1, bx1)
    oy = min(ay2, by2) - max(ay1, by1)
    if ox > TOL and oy > TOL:
        return round(ox, 3), round(oy, 3)
    return None


problems = 0
for i, sl in enumerate(prs.slides, 1):
    shapes = list(sl.shapes)
    # 1. di luar kanvas
    for sh in shapes:
        if sh.left is None:
            continue
        l, t = sh.left / EMU, sh.top / EMU
        r, b = l + sh.width / EMU, t + sh.height / EMU
        if l < -0.01 or t < -0.01 or r > SW + 0.01 or b > SH + 0.01:
            print(f"[{i}] LUAR KANVAS  {sh.shape_type}  l={l:.2f} t={t:.2f} "
                  f"r={r:.2f} b={b:.2f}  teks={sh.has_text_frame and sh.text_frame.text[:30]!r}")
            problems += 1
    # 2. teks bertumpuk (kotak teks lawan kotak teks)
    ts = text_shapes(sl)
    for a in range(len(ts)):
        for b in range(a + 1, len(ts)):
            # kotak teks di dalam shape berisi (callout/kartu) sengaja tidak dicek
            # karena keduanya adalah shape yang sama.
            ov = rects_overlap(ts[a], ts[b])
            if ov:
                print(f"[{i}] TUMPANG TINDIH TEKS  {ts[a].text_frame.text[:26]!r} "
                      f"x {ts[b].text_frame.text[:26]!r}  ({ov[0]}\" x {ov[1]}\")")
                problems += 1
    # 3. luapan teks
    for sh in ts:
        have = sh.height / EMU
        want = need_height(sh)
        if want > have + 0.05:
            print(f"[{i}] LUAPAN  butuh {want:.2f}\" tersedia {have:.2f}\"  "
                  f"{sh.text_frame.text[:44]!r}")
            problems += 1

print("---")
print("slide:", len(prs.slides._sldIdLst), "| temuan:", problems)
